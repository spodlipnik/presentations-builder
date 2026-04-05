"""Pytest fixtures — generates real PPTX files for integration testing."""
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


@pytest.fixture
def simple_title_pptx(tmp_path):
    """A minimal 1-slide PPTX with a title text box (for shape/font tests)."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    tf_shape = slide.shapes.add_textbox(
        left=Inches(1.0),
        top=Inches(2.0),
        width=Inches(11.0),
        height=Inches(1.5),
    )
    tf = tf_shape.text_frame
    tf.text = "Sample Title"
    run = tf.paragraphs[0].runs[0]
    run.font.name = "Inter"
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)

    out = tmp_path / "simple_title.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def multi_slide_pptx(tmp_path):
    """A 3-slide PPTX: title, assertion-evidence-like, and closing."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: title-like (single large centered text)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s1.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.0), Inches(1.5))
    tb.text_frame.text = "Microbiota Cutánea"
    tb.text_frame.paragraphs[0].runs[0].font.size = Pt(54)
    tb.text_frame.paragraphs[0].runs[0].font.bold = True

    # Slide 2: assertion-evidence-like (headline + image box)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    h = s2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.0), Inches(1.0))
    h.text_frame.text = "La incidencia aumentó 5% anual"
    h.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
    h.text_frame.paragraphs[0].runs[0].font.bold = True
    # Image placeholder (rectangle standing in for image)
    s2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5))

    # Slide 3: closing (small text)
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    c = s3.shapes.add_textbox(Inches(4.0), Inches(3.0), Inches(5.0), Inches(1.0))
    c.text_frame.text = "Preguntas?"
    c.text_frame.paragraphs[0].runs[0].font.size = Pt(48)

    out = tmp_path / "multi_slide.pptx"
    prs.save(str(out))
    return out


@pytest.fixture
def grouped_shapes_pptx(tmp_path):
    """PPTX with a group shape containing nested text boxes."""
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Add two text boxes (python-pptx can't easily create groups post-hoc,
    # so we simulate with separate shapes for now; extractor should handle both)
    slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(3.0), Inches(1.0)).text_frame.text = "Grouped A"
    slide.shapes.add_textbox(Inches(5.0), Inches(1.0), Inches(3.0), Inches(1.0)).text_frame.text = "Grouped B"

    out = tmp_path / "grouped.pptx"
    prs.save(str(out))
    return out
