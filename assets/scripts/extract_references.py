#!/usr/bin/env python3
"""Extract structured layout catalog from a reference PPTX file.

Usage:
    python3 extract_references.py --input ref.pptx --output catalog.yaml --theme-name my-theme
"""
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
from collections import Counter


EMU_PER_INCH = 914400


def load_presentation(pptx_path):
    """Load a PPTX file and return the Presentation object."""
    return Presentation(str(pptx_path))


def get_slide_dimensions(prs):
    """Return slide dimensions as dict with width_in, height_in, aspect_ratio."""
    width_in = round(prs.slide_width / EMU_PER_INCH, 3)
    height_in = round(prs.slide_height / EMU_PER_INCH, 3)
    ratio = width_in / height_in
    # Common ratios
    if abs(ratio - 16/9) < 0.01:
        aspect = "16:9"
    elif abs(ratio - 4/3) < 0.01:
        aspect = "4:3"
    else:
        aspect = f"{width_in}:{height_in}"
    return {
        "width_in": width_in,
        "height_in": height_in,
        "aspect_ratio": aspect,
    }


def extract_shape_box(shape, slide_width_emu, slide_height_emu):
    """Return [x, y, w, h] as fractions (0-1) of slide dimensions.

    Fractions are resolution-independent and readable.
    """
    return [
        round(shape.left / slide_width_emu, 3),
        round(shape.top / slide_height_emu, 3),
        round(shape.width / slide_width_emu, 3),
        round(shape.height / slide_height_emu, 3),
    ]


# OOXML namespace for drawingml (used by font inheritance parsers)
_A_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _parse_rpr_element(rpr_elem, theme_fonts=None, theme_colors=None):
    """Parse a <a:rPr> or <a:defRPr> XML element into a partial font dict.

    Resolves theme font references (+mn-lt, +mj-lt, etc.) and schemeClr colors.
    Returns dict with any of: family, size_pt, weight, color_rgb.
    """
    if rpr_elem is None:
        return {}
    result = {}
    # Size: sz attribute in 100ths of point
    sz = rpr_elem.get("sz")
    if sz:
        try:
            result["size_pt"] = int(sz) / 100.0
        except ValueError:
            pass
    # Weight: b="1" bold
    b = rpr_elem.get("b")
    if b == "1":
        result["weight"] = 700
    elif b == "0":
        result["weight"] = 400
    # Family: check latin first, then sym (Keynote uses sym for concrete font)
    family = None
    latin = rpr_elem.find("a:latin", _A_NS)
    if latin is not None:
        tf_val = latin.get("typeface", "")
        if tf_val.startswith("+") and theme_fonts:
            # Theme font reference: +mn-lt, +mj-lt, +mn-ea, +mj-ea, +mn-cs, +mj-cs
            token = tf_val[1:]  # strip leading +
            key_map = {
                "mn-lt": "minor_latin", "mj-lt": "major_latin",
                "mn-ea": "minor_ea", "mj-ea": "major_ea",
                "mn-cs": "minor_cs", "mj-cs": "major_cs",
            }
            family = theme_fonts.get(key_map.get(token))
        elif tf_val:
            family = tf_val
    # sym typeface often overrides (Keynote puts concrete font here)
    sym = rpr_elem.find("a:sym", _A_NS)
    if sym is not None:
        sym_tf = sym.get("typeface", "")
        if sym_tf and not sym_tf.startswith("+"):
            family = sym_tf
    if family:
        result["family"] = family
    # Color: look inside solidFill child
    solid = rpr_elem.find("a:solidFill", _A_NS)
    if solid is not None:
        srgb = solid.find("a:srgbClr", _A_NS)
        scheme = solid.find("a:schemeClr", _A_NS)
        if srgb is not None:
            val = srgb.get("val", "").upper()
            if val:
                result["color_rgb"] = val
        elif scheme is not None and theme_colors:
            role = scheme.get("val")
            if role and role in theme_colors:
                result["color_rgb"] = theme_colors[role]
    return result


def extract_font_info(shape, theme_fonts=None, theme_colors=None):
    """Extract font info from first run of first paragraph of a text shape.

    Walks inheritance chain when run-level fields are empty (Keynote exports):
      run.rPr → paragraph.pPr/defRPr → text_frame.lstStyle/lvl1pPr/defRPr

    Returns None if shape has no text frame or no runs.
    Returns dict with family, size_pt, weight, color_rgb (may have None values).
    """
    if not getattr(shape, "has_text_frame", False):
        return None
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        return None

    run = tf.paragraphs[0].runs[0]
    font = run.font

    # Weight from python-pptx run level
    if font.bold is True:
        weight = 700
    elif font.bold is False:
        weight = 400
    else:
        weight = None

    # RGB color at run level
    color_rgb = None
    try:
        if font.color.type is not None and hasattr(font.color, "rgb") and font.color.rgb is not None:
            color_rgb = str(font.color.rgb)
    except (AttributeError, TypeError):
        pass

    family = font.name
    size_pt = font.size.pt if font.size else None

    # If any field is missing, walk XML inheritance chain
    if None in (family, size_pt, weight, color_rgb):
        inherited = {}
        # 1. run's own rPr (in case python-pptx missed something)
        rpr = run._r.find("a:rPr", _A_NS)
        if rpr is not None:
            inherited.update(_parse_rpr_element(rpr, theme_fonts, theme_colors))
        # 2. paragraph's defRPr inside pPr
        if tf.paragraphs[0]._pPr is not None:
            p_defrpr = tf.paragraphs[0]._pPr.find("a:defRPr", _A_NS)
            if p_defrpr is not None:
                for k, v in _parse_rpr_element(p_defrpr, theme_fonts, theme_colors).items():
                    inherited.setdefault(k, v)
        # 3. text_frame's lstStyle (Keynote style)
        lst_style = tf._txBody.find("a:lstStyle", _A_NS)
        if lst_style is not None:
            # Pick paragraph level (lvl1pPr for level 0, etc.)
            para_level = tf.paragraphs[0].level or 0
            lvl_tag = f"a:lvl{para_level + 1}pPr"
            lvl_elem = lst_style.find(lvl_tag, _A_NS)
            if lvl_elem is None:
                lvl_elem = lst_style.find("a:lvl1pPr", _A_NS)  # fallback
            if lvl_elem is not None:
                defrpr = lvl_elem.find("a:defRPr", _A_NS)
                if defrpr is not None:
                    for k, v in _parse_rpr_element(defrpr, theme_fonts, theme_colors).items():
                        inherited.setdefault(k, v)

        # Apply inherited values where run-level was None
        family = family or inherited.get("family")
        size_pt = size_pt or inherited.get("size_pt")
        weight = weight or inherited.get("weight")
        color_rgb = color_rgb or inherited.get("color_rgb")

    return {
        "family": family,
        "size_pt": size_pt,
        "weight": weight,
        "color_rgb": color_rgb,
    }


# OOXML namespaces for theme XML
_THEME_NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}

# Order of color roles in theme1.xml clrScheme element
_THEME_COLOR_ORDER = [
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
]


def extract_theme_colors(pptx_path):
    """Parse ppt/theme/theme1.xml from PPTX zip and return role→hex color map.

    Returns dict like {"dk1": "000000", "lt1": "FFFFFF", "accent1": "4472C4", ...}
    """
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        # Find first theme file (most PPTX have exactly one)
        theme_files = [n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
        if not theme_files:
            return {}
        with z.open(theme_files[0]) as f:
            tree = etree.parse(f)

    root = tree.getroot()
    # Find <a:clrScheme> which contains the 12 theme colors in order
    clr_scheme = root.find(".//a:clrScheme", _THEME_NS)
    if clr_scheme is None:
        return {}

    colors = {}
    # Each child of clrScheme is a color role: <a:dk1>, <a:lt1>, etc.
    # Inside each is either <a:srgbClr val="HEX"/> or <a:sysClr lastClr="HEX"/>
    for role_name in _THEME_COLOR_ORDER:
        role_elem = clr_scheme.find(f"a:{role_name}", _THEME_NS)
        if role_elem is None:
            continue
        srgb = role_elem.find("a:srgbClr", _THEME_NS)
        sysclr = role_elem.find("a:sysClr", _THEME_NS)
        if srgb is not None:
            colors[role_name] = srgb.get("val", "").upper()
        elif sysclr is not None:
            colors[role_name] = sysclr.get("lastClr", "").upper()

    return colors


def extract_theme_fonts(pptx_path):
    """Parse theme1.xml fontScheme to return major/minor font families.

    Returns dict like {"major_latin": "Calibri Light", "minor_latin": "Calibri", ...}
    Keys: major_latin, minor_latin, major_ea, minor_ea, major_cs, minor_cs (any may be absent).
    """
    with zipfile.ZipFile(str(pptx_path), "r") as z:
        theme_files = [n for n in z.namelist() if n.startswith("ppt/theme/theme") and n.endswith(".xml")]
        if not theme_files:
            return {}
        with z.open(theme_files[0]) as f:
            tree = etree.parse(f)

    font_scheme = tree.getroot().find(".//a:fontScheme", _THEME_NS)
    if font_scheme is None:
        return {}

    result = {}
    for role, elem_name in [("major", "majorFont"), ("minor", "minorFont")]:
        parent = font_scheme.find(f"a:{elem_name}", _THEME_NS)
        if parent is None:
            continue
        for script in ("latin", "ea", "cs"):
            tag = parent.find(f"a:{script}", _THEME_NS)
            if tag is not None:
                tf_val = tag.get("typeface", "")
                if tf_val:
                    result[f"{role}_{script}"] = tf_val
    return result


def extract_shape_type(shape):
    """Categorize shape as: text_box, image, table, chart, group, placeholder, or other.

    Handles SmartArt safely (returns 'other' without raising).
    """
    try:
        st = shape.shape_type
    except (AttributeError, NotImplementedError):
        return "other"

    # Check graphic frames first (tables/charts live inside these)
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "placeholder"
    if getattr(shape, "has_text_frame", False):
        return "text_box"
    return "other"


def iter_shapes_recursive(shapes):
    """Yield all shapes, recursing into groups. Skips SmartArt safely.

    Returns (shape, z_order) tuples. z_order is the iteration index (0-based,
    first = back).
    """
    z = 0
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for inner_shape, _ in iter_shapes_recursive(shape.shapes):
                    yield inner_shape, z
                    z += 1
                continue
        except (AttributeError, NotImplementedError):
            # SmartArt or unsupported — skip gracefully
            continue
        yield shape, z
        z += 1


def extract_slide(slide, slide_number, slide_width_emu, slide_height_emu,
                  theme_fonts=None, theme_colors=None):
    """Extract complete structured info for a single slide.

    Returns dict with slide_number, shapes (list), inferred_role (None initially).
    theme_fonts and theme_colors are passed to font extraction for inheritance resolution.
    """
    shapes_info = []
    for shape, z_order in iter_shapes_recursive(slide.shapes):
        shape_dict = {
            "type": extract_shape_type(shape),
            "box": extract_shape_box(shape, slide_width_emu, slide_height_emu),
            "z_order": z_order,
        }
        # Text content + font (if applicable)
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text or ""
            shape_dict["text_sample"] = text[:200]  # cap at 200 chars
            shape_dict["font"] = extract_font_info(shape, theme_fonts, theme_colors)
            # Paragraph alignment
            if shape.text_frame.paragraphs:
                pa = shape.text_frame.paragraphs[0].alignment
                shape_dict["alignment"] = str(pa).split(".")[-1].lower() if pa else None
        # Placeholder info
        if getattr(shape, "is_placeholder", False):
            shape_dict["is_placeholder"] = True
            try:
                shape_dict["placeholder_type"] = str(shape.placeholder_format.type).split(".")[-1].lower()
            except AttributeError:
                shape_dict["placeholder_type"] = None
        shapes_info.append(shape_dict)

    return {
        "slide_number": slide_number,
        "shapes": shapes_info,
        "inferred_role": None,  # filled in by role inference step
    }


def detect_dominant_tokens(slides_data):
    """Count font families and colors across all shapes to identify dominant tokens.

    Returns dict with dominant_fonts and dominant_colors (sorted by count, descending).
    """
    font_counter = Counter()
    color_counter = Counter()
    weight_per_family = {}  # family → Counter of weights

    for slide in slides_data:
        for shape in slide["shapes"]:
            font = shape.get("font")
            if not font:
                continue
            family = font.get("family")
            if family:
                font_counter[family] += 1
                weight_per_family.setdefault(family, Counter())
                if font.get("weight"):
                    weight_per_family[family][font["weight"]] += 1
            color = font.get("color_rgb")
            if color:
                color_counter[color] += 1

    dominant_fonts = [
        {
            "family": family,
            "count": count,
            "weight_distribution": dict(weight_per_family.get(family, {})),
        }
        for family, count in font_counter.most_common(5)
    ]
    dominant_colors = [
        {"hex": f"#{hex_str}", "count": count}
        for hex_str, count in color_counter.most_common(10)
    ]

    return {
        "dominant_fonts": dominant_fonts,
        "dominant_colors": dominant_colors,
    }


def infer_role(slide_info):
    """Apply heuristics to guess which canonical role a slide represents.

    Returns one of the 18 role IDs, or 'unknown' if no rule matches.
    User validates/corrects in talk-theme-builder Fase 2.
    """
    shapes = slide_info.get("shapes", [])
    if not shapes:
        return "unknown"

    # Classify shapes
    text_shapes = [s for s in shapes if s["type"] == "text_box" and s.get("text_sample")]
    image_or_empty_boxes = [s for s in shapes if s["type"] in ("image", "placeholder")
                            or (s["type"] == "text_box" and not s.get("text_sample"))]

    # Helper: large text if font size ≥ 36pt
    def is_large(shape):
        f = shape.get("font") or {}
        return (f.get("size_pt") or 0) >= 36

    # Title: single large text, centered, no other content
    if len(text_shapes) == 1 and not image_or_empty_boxes and is_large(text_shapes[0]):
        return "title"
    if len(shapes) == 1 and text_shapes and is_large(text_shapes[0]):
        return "title"

    # Quote: large centered text, contains quote marks
    if len(text_shapes) == 1 and is_large(text_shapes[0]):
        txt = text_shapes[0].get("text_sample", "")
        if '"' in txt or '"' in txt or '«' in txt:
            return "quote-pullout"

    # Section divider: 1-2 text shapes, very large font, minimal content
    if len(shapes) <= 2 and text_shapes and any(((s.get("font") or {}).get("size_pt") or 0) >= 60 for s in text_shapes):
        return "section-divider"

    # Assertion-evidence: 1 small/medium headline text at top + content below
    top_shapes = [s for s in shapes if s["box"][1] < 0.25]
    bottom_shapes = [s for s in shapes if s["box"][1] >= 0.25]
    if len(top_shapes) == 1 and top_shapes[0]["type"] == "text_box" and bottom_shapes:
        return "assertion-evidence"

    # Image gallery: multiple image/empty shapes
    if len([s for s in shapes if s["type"] in ("image", "placeholder")]) >= 3:
        return "image-gallery"

    # Image-fullbleed: single large image covering most of slide
    large_images = [s for s in shapes if s["type"] == "image"
                    and s["box"][2] >= 0.8 and s["box"][3] >= 0.7]
    if large_images:
        return "image-fullbleed"

    # Comparison: 2 shapes roughly side by side
    if len(shapes) == 2:
        s1, s2 = shapes[0], shapes[1]
        if abs(s1["box"][1] - s2["box"][1]) < 0.1 and abs(s1["box"][3] - s2["box"][3]) < 0.1:
            return "comparison"

    return "unknown"


import yaml
import argparse
from datetime import datetime


def build_catalog(pptx_path, theme_name=None):
    """Build complete catalog dict from a PPTX file.

    This is the main orchestration function that combines all extraction steps.
    """
    pptx_path = Path(pptx_path)
    prs = load_presentation(pptx_path)

    # Theme info loaded once for inheritance resolution during per-slide extraction
    theme_colors = extract_theme_colors(pptx_path)
    theme_fonts = extract_theme_fonts(pptx_path)

    # Extract per-slide data (pass theme info for font inheritance)
    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_info = extract_slide(
            slide, i + 1, prs.slide_width, prs.slide_height,
            theme_fonts=theme_fonts, theme_colors=theme_colors,
        )
        slide_info["inferred_role"] = infer_role(slide_info)
        slides_data.append(slide_info)

    # Detect global tokens (now that fonts are resolved via inheritance)
    tokens = detect_dominant_tokens(slides_data)
    dimensions = get_slide_dimensions(prs)

    # Get python-pptx version
    try:
        import pptx
        pptx_version = pptx.__version__
    except AttributeError:
        pptx_version = "unknown"

    return {
        "meta": {
            "source_file": pptx_path.name,
            "extracted_at": datetime.utcnow().isoformat() + "Z",
            "slides_count": len(slides_data),
            "python_pptx_version": pptx_version,
            "theme_name": theme_name,
        },
        "theme_detection": {
            "dominant_colors": tokens["dominant_colors"],
            "dominant_fonts": tokens["dominant_fonts"],
            "theme_xml_colors": theme_colors,
            "theme_xml_fonts": theme_fonts,
            "slide_dimensions": dimensions,
        },
        "slides": slides_data,
    }


def write_catalog_yaml(catalog, output_path):
    """Write catalog dict to YAML file."""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "w", encoding="utf-8") as f:
        yaml.safe_dump(catalog, f, sort_keys=False, allow_unicode=True, default_flow_style=False)


def main():
    parser = argparse.ArgumentParser(
        description="Extract structured layout catalog from a PPTX reference file."
    )
    parser.add_argument("--input", "-i", required=True, help="Input PPTX file path")
    parser.add_argument("--output", "-o", required=True, help="Output YAML file path")
    parser.add_argument("--theme-name", help="Optional theme name to include in metadata")
    args = parser.parse_args()

    catalog = build_catalog(args.input, theme_name=args.theme_name)
    write_catalog_yaml(catalog, args.output)
    print(f"Extracted {catalog['meta']['slides_count']} slides \u2192 {args.output}")


if __name__ == "__main__":
    main()
